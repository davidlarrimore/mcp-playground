import logging
import os
import sys
from pathlib import Path
from typing import List, Optional
from fastmcp import FastMCP
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Add parent directory to path for shared module
sys.path.insert(0, str(Path(__file__).parent.parent.parent))
from shared.download_urls import generate_signed_url

LOG_LEVEL = os.getenv("LOG_LEVEL", "INFO").upper()
logging.basicConfig(level=LOG_LEVEL, format="%(asctime)s [%(levelname)s] %(name)s - %(message)s")
logger = logging.getLogger("analytics-mcp")

WORKSPACE = Path(os.getenv("WORKSPACE", "/workspace")).resolve()

mcp = FastMCP("analytics")

@mcp.tool
async def save_uploaded_file(
    filename: str,
    content_base64: str
) -> dict:
    """
    Save an uploaded file (provided as base64) to the workspace for analysis.

    This is a convenience tool for when users upload files via chat.
    After saving, the file will be accessible to all analytics tools.

    IMPORTANT: Preserve the original file extension! If user uploads "data.xlsx",
    save it as "data.xlsx" (or similar .xlsx name), NOT as "data.txt".

    Args:
        filename: Name to save the file as - MUST keep original file extension (e.g., "data.xlsx", "report.csv", "document.pdf")
        content_base64: Base64-encoded file content from the uploaded file
    """
    import base64

    try:
        # Decode base64 content
        file_data = base64.b64decode(content_base64)

        # Save to workspace
        output_path = WORKSPACE / filename
        output_path.write_bytes(file_data)

        logger.info(f"Saved uploaded file: {output_path}")

        return {
            "path": str(output_path.relative_to(WORKSPACE)),
            "absolute_path": str(output_path),
            "size": output_path.stat().st_size,
            "message": f"File saved successfully. Use '{filename}' (without 'docs/' prefix) in other analytics tools."
        }
    except Exception as e:
        error_msg = f"Error saving file: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

@mcp.tool
async def list_data_files(
    pattern: Optional[str] = None
) -> dict:
    """
    List available data files in the workspace.

    IMPORTANT: Files uploaded via chat are NOT automatically available in the workspace.
    To work with uploaded files, users must first use the filesystem MCP to save them.

    Args:
        pattern: Optional glob pattern to filter files (e.g., "*.csv", "border_*.xlsx")
    """
    if pattern:
        files = list(WORKSPACE.glob(pattern))
    else:
        files = [f for f in WORKSPACE.iterdir() if f.is_file() and f.suffix in ['.csv', '.xlsx', '.xls']]

    file_list = []
    for f in sorted(files):
        file_list.append({
            "name": f.name,
            "path": str(f.relative_to(WORKSPACE)),
            "size": f.stat().st_size,
            "type": f.suffix
        })

    logger.info(f"Found {len(file_list)} data files")

    return {
        "files": file_list,
        "count": len(file_list),
        "workspace": str(WORKSPACE),
        "note": "Files uploaded via chat must be saved to workspace using filesystem MCP before analytics tools can access them"
    }

@mcp.tool
async def merge_excel_files(
    file_paths: List[str],
    output_filename: str
) -> dict:
    """
    Merge multiple data files (CSV or Excel) into a single dataset.

    Args:
        file_paths: List of file paths (relative to workspace). Supports CSV and Excel files. Do NOT include 'docs/' prefix.
        output_filename: Output filename for merged data (CSV or Excel)
    """
    dfs = []
    missing_files = []

    for fp in file_paths:
        full_path = WORKSPACE / fp

        # Check if file exists
        if not full_path.exists():
            missing_files.append(fp)
            continue

        # Auto-detect file type and read accordingly
        try:
            if fp.endswith('.csv'):
                df = pd.read_csv(full_path)
            elif fp.endswith(('.xlsx', '.xls')):
                df = pd.read_excel(full_path)
            else:
                # Try CSV first, then Excel as fallback
                try:
                    df = pd.read_csv(full_path)
                except:
                    df = pd.read_excel(full_path)

            dfs.append(df)
            logger.info(f"Successfully read {fp}: {len(df)} rows")
        except Exception as e:
            logger.error(f"Error reading {fp}: {str(e)}")
            missing_files.append(f"{fp} (error: {str(e)})")

    if missing_files:
        error_msg = f"Could not read files: {', '.join(missing_files)}"
        logger.error(error_msg)
        return {
            "error": error_msg,
            "missing_files": missing_files,
            "successfully_read": len(dfs)
        }

    if not dfs:
        return {"error": "No files were successfully read"}

    merged = pd.concat(dfs, ignore_index=True)

    output_path = WORKSPACE / output_filename
    if output_filename.endswith('.csv'):
        merged.to_csv(output_path, index=False)
    else:
        merged.to_excel(output_path, index=False)

    logger.info(f"Merged {len(file_paths)} files into {output_path}")

    return {
        "path": str(output_path.relative_to(WORKSPACE)),
        "absolute_path": str(output_path),
        "rows": len(merged),
        "columns": list(merged.columns),
        "source_files": len(dfs)
    }

@mcp.tool
async def calculate_summary_stats(
    file_path: str,
    group_by: Optional[str] = None
) -> dict:
    """
    Calculate summary statistics from data.

    Args:
        file_path: Path to data file (CSV or Excel). Do NOT include 'docs/' prefix.
        group_by: Optional column to group by (e.g., "Region")
    """
    full_path = WORKSPACE / file_path

    # Check if file exists
    if not full_path.exists():
        # Try to be helpful about common mistakes
        alt_path = WORKSPACE / file_path.replace("docs/", "")
        if alt_path.exists():
            suggestion = f"Did you mean '{file_path.replace('docs/', '')}' instead? (Don't include 'docs/' prefix)"
        else:
            suggestion = "Use list_data_files tool to see available files. Files uploaded via chat must first be saved using filesystem MCP."

        error_msg = f"File not found: {file_path}. {suggestion}"
        logger.error(error_msg)
        return {"error": error_msg}

    try:
        if file_path.endswith('.csv'):
            df = pd.read_csv(full_path)
        else:
            df = pd.read_excel(full_path)
    except Exception as e:
        error_msg = f"Error reading file {file_path}: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

    if group_by and group_by in df.columns:
        summary = df.groupby(group_by).agg({
            col: ['sum', 'mean', 'count']
            for col in df.select_dtypes(include='number').columns
        }).to_dict()
    else:
        summary = df.describe().to_dict()

    logger.info(f"Calculated stats for {file_path}")

    return {
        "summary": summary,
        "total_rows": len(df),
        "columns": list(df.columns)
    }

@mcp.tool
async def generate_chart(
    file_path: str,
    chart_type: str,
    x_column: str,
    y_column: str,
    output_filename: str,
    title: Optional[str] = None
) -> dict:
    """
    Generate a chart from data.

    Args:
        file_path: Path to data file (CSV or Excel). Do NOT include 'docs/' prefix.
        chart_type: Type of chart ('bar', 'line', 'pie')
        x_column: Column for x-axis
        y_column: Column for y-axis
        output_filename: Output image filename (PNG)
        title: Optional chart title
    """
    full_path = WORKSPACE / file_path

    # Check if file exists
    if not full_path.exists():
        # Try to be helpful about common mistakes
        alt_path = WORKSPACE / file_path.replace("docs/", "")
        if alt_path.exists():
            suggestion = f"Did you mean '{file_path.replace('docs/', '')}' instead? (Don't include 'docs/' prefix)"
        else:
            suggestion = "Use list_data_files tool to see available files. Files uploaded via chat must first be saved using filesystem MCP."

        error_msg = f"File not found: {file_path}. {suggestion}"
        logger.error(error_msg)
        return {"error": error_msg}

    try:
        if file_path.endswith('.csv'):
            df = pd.read_csv(full_path)
        else:
            df = pd.read_excel(full_path)
    except Exception as e:
        error_msg = f"Error reading file {file_path}: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

    # Validate columns exist
    if x_column not in df.columns:
        return {"error": f"Column '{x_column}' not found in data. Available columns: {list(df.columns)}"}
    if y_column not in df.columns:
        return {"error": f"Column '{y_column}' not found in data. Available columns: {list(df.columns)}"}

    try:
        plt.figure(figsize=(10, 6))

        if chart_type == 'bar':
            plt.bar(df[x_column], df[y_column])
        elif chart_type == 'line':
            plt.plot(df[x_column], df[y_column], marker='o')
        elif chart_type == 'pie':
            plt.pie(df[y_column], labels=df[x_column], autopct='%1.1f%%')
        else:
            return {"error": f"Unsupported chart type: {chart_type}. Use 'bar', 'line', or 'pie'."}

        if title:
            plt.title(title)
        if chart_type != 'pie':
            plt.xlabel(x_column)
            plt.ylabel(y_column)
        plt.tight_layout()

        output_path = WORKSPACE / output_filename
        plt.savefig(output_path, dpi=150, bbox_inches='tight')
        plt.close()

        logger.info(f"Generated {chart_type} chart: {output_path}")

        # Generate signed download URL
        try:
            download_info = generate_signed_url(str(output_path))
            download_url = download_info["url"]
            expires_hours = download_info["expires_in"] // 3600
        except Exception as e:
            logger.error(f"Error generating download URL: {e}")
            download_info = None
            download_url = None
            expires_hours = 0

        result = {
            "path": str(output_path.relative_to(WORKSPACE)),
            "absolute_path": str(output_path),
            "size": output_path.stat().st_size,
            "chart_type": chart_type,
            "message": f"{chart_type.capitalize()} chart created successfully: {title or output_filename}"
        }

        if download_info:
            result["download_url"] = download_url
            result["download_expires_at"] = download_info["expires_at"]
            result["download_expires_in"] = download_info["expires_in"]
            result["message"] += f"\n\n📥 Download your chart:\n{download_url}\n\n⏰ Link expires in {expires_hours} hours"

        return result
    except Exception as e:
        plt.close()  # Clean up on error
        error_msg = f"Error generating chart: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

@mcp.tool
async def create_presentation(
    title: str,
    slides: List[dict],
    output_filename: str,
    subtitle: Optional[str] = None
) -> dict:
    """
    Create a PowerPoint presentation.

    Args:
        title: Presentation title
        slides: List of slides with structure:
            [{
                "title": "Slide Title",
                "content": "Slide content text",
                "bullet_points": ["Point 1", "Point 2"],
                "image_path": "optional/path/to/image.png",
                "table": {
                    "columns": ["Region", "Revenue"],
                    "rows": [{"Region": "East", "Revenue": "50K"}]
                }
            }]
        output_filename: Name of output file
        subtitle: Optional subtitle for title slide
    """
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle

        # Content slides
        for slide_def in slides:
            # Use title and content layout
            content_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_layout)

            # Add title
            slide.shapes.title.text = slide_def.get("title", "")

            # Get content placeholder
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]

                # Add bullet points if provided
                if "bullet_points" in slide_def:
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    for point in slide_def["bullet_points"]:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0
                # Add content text if provided
                elif "content" in slide_def:
                    content_placeholder.text = slide_def["content"]

            # Add image if provided
            if "image_path" in slide_def:
                image_full_path = WORKSPACE / slide_def["image_path"]
                if image_full_path.exists():
                    left = Inches(5)
                    top = Inches(2)
                    slide.shapes.add_picture(
                        str(image_full_path),
                        left, top,
                        height=Inches(4)
                    )

            # Add table if provided
            if "table" in slide_def:
                table_data = slide_def["table"]
                rows = len(table_data["rows"]) + 1  # +1 for header
                cols = len(table_data["columns"])

                left = Inches(1)
                top = Inches(2.5)
                width = Inches(8)
                height = Inches(0.5) * rows

                table = slide.shapes.add_table(rows, cols, left, top, width, height).table

                # Header row
                for col_idx, col_name in enumerate(table_data["columns"]):
                    cell = table.cell(0, col_idx)
                    cell.text = col_name
                    cell.text_frame.paragraphs[0].font.bold = True

                # Data rows
                for row_idx, row_data in enumerate(table_data["rows"], 1):
                    for col_idx, col_name in enumerate(table_data["columns"]):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(row_data.get(col_name, ""))

        output_path = WORKSPACE / output_filename
        prs.save(output_path)

        logger.info(f"Created PowerPoint presentation: {output_path}")

        # Generate signed download URL
        try:
            download_info = generate_signed_url(str(output_path))
            download_url = download_info["url"]
            expires_hours = download_info["expires_in"] // 3600
        except Exception as e:
            logger.error(f"Error generating download URL: {e}")
            download_info = None
            download_url = None
            expires_hours = 0

        total_slides = len(slides) + 1  # +1 for title slide
        result = {
            "path": str(output_path.relative_to(WORKSPACE)),
            "absolute_path": str(output_path),
            "size": output_path.stat().st_size,
            "slides": total_slides,
            "message": f"PowerPoint presentation '{title}' created successfully with {total_slides} slide(s)."
        }

        if download_info:
            result["download_url"] = download_url
            result["download_expires_at"] = download_info["expires_at"]
            result["download_expires_in"] = download_info["expires_in"]
            result["message"] += f"\n\n📥 Download your presentation:\n{download_url}\n\n⏰ Link expires in {expires_hours} hours"

        return result

    except Exception as e:
        error_msg = f"Error creating presentation: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

if __name__ == "__main__":
    mcp.run()
